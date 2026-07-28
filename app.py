"""
Agente Corporativo de IA (RAG) - Desafío Alura Agentes
-------------------------------------------------------------------
Agente conversacional que responde preguntas de los colaboradores de
una empresa con base UNICAMENTE en los documentos internos cargados
(RH, Financiero, Legal, Operacional, Estratégico, Marketing, Datos y
Sistemas, I+D, Calidad, Comunicación Interna, etc). Acceso abierto a
cualquier colaborador, sin restricciones.

Formatos soportados: PDF, Word (.docx), Excel (.xlsx/.xls), PowerPoint
(.pptx), Markdown (.md), CSV, JSON y HTML.

Motor de lenguaje: API de Google Gemini (capa gratuita, sin necesidad
de tarjeta de crédito). Se eligió una API en la nube (en vez de un
modelo local con Ollama) porque el despliegue se hace en Streamlit
Community Cloud, que no permite instalar ni correr servicios como
Ollama.

Requisitos previos:
    1. Obtener una API key gratuita en Google AI Studio
       (https://aistudio.google.com/apikey)
    2. Instalar dependencias:
         pip install -r requirements.txt
    3. Configurar la API key (ver README.md) como variable de entorno
       GEMINI_API_KEY o en .streamlit/secrets.toml
    4. Ejecutar:
         streamlit run app.py
"""

import os
import io
import json
import numpy as np
import pandas as pd
import streamlit as st
import faiss
from google import genai
from google.genai import types as genai_types
from pypdf import PdfReader
import docx
from pptx import Presentation
from bs4 import BeautifulSoup

# --------------------------------------------------------------------------
# Configuración general
# --------------------------------------------------------------------------
st.set_page_config(page_title="Agente Corporativo IA", page_icon="🏢", layout="wide")

CHUNK_SIZE = 900          # caracteres por fragmento
CHUNK_OVERLAP = 150       # solapamiento entre fragmentos
TOP_K = 4                 # cantidad de fragmentos relevantes a recuperar por defecto
EMBED_MODEL_NAME = "paraphrase-multilingual-MiniLM-L12-v2"  # embeddings locales, soporta español

MODELOS_GEMINI = ["gemini-2.5-flash", "gemini-3.1-flash-lite", "gemini-3-flash"]

CATEGORIAS_REFERENCIA = [
    "Recursos Humanos (políticas, beneficios, onboarding)",
    "Financiero y Contable (estados de resultados, balances, políticas de gastos)",
    "Operacional (procesos, procedimientos, manuales técnicos)",
    "Estratégico (planes, roadmaps)",
    "Legal y Compliance (contratos, normativas de privacidad/protección de datos)",
    "Marketing y Comercial (pitch decks, tablas de precios)",
    "Datos y Sistemas (planillas, bases de clientes, APIs)",
    "Investigación y Desarrollo (market research, casos de negocio)",
    "Calidad (auditorías, planes correctivos)",
    "Comunicación Interna (comunicados, minutas, newsletters)",
]

SYSTEM_PROMPT = (
    "Eres el Agente Corporativo de IA de la empresa. Respondes preguntas de "
    "cualquier colaborador (acceso abierto, sin restricciones) sobre "
    "documentos internos que abarcan distintas áreas: Recursos Humanos, "
    "Financiero/Contable, Operacional, Estratégico, Legal/Compliance, "
    "Marketing/Comercial, Datos y Sistemas, Investigación y Desarrollo, "
    "Calidad y Comunicación Interna.\n\n"
    "Reglas estrictas:\n"
    "1. Responde EXCLUSIVAMENTE con base en el CONTEXTO proporcionado, "
    "extraído de los documentos cargados. No uses conocimiento externo ni "
    "inventes datos.\n"
    "2. Si la respuesta no se encuentra en el CONTEXTO, responde exactamente: "
    "'No encuentro esa información en los documentos cargados.'\n"
    "3. Cuando respondas, sé claro y conciso, y cita el documento de origen "
    "entre corchetes, ej. [politica_vacaciones.pdf].\n"
    "4. Responde siempre en español, con un tono profesional y corporativo."
)

# --------------------------------------------------------------------------
# Utilidades de extracción de texto por formato
# --------------------------------------------------------------------------
def extraer_texto_pdf(archivo) -> str:
    lector = PdfReader(archivo)
    return "\n".join(pagina.extract_text() or "" for pagina in lector.pages)


def extraer_texto_docx(archivo) -> str:
    documento = docx.Document(archivo)
    return "\n".join(parrafo.text for parrafo in documento.paragraphs)


def extraer_texto_txt(archivo) -> str:
    return archivo.read().decode("utf-8", errors="ignore")


def _dataframe_a_texto(df: pd.DataFrame, encabezado: str = "") -> str:
    lineas = [encabezado] if encabezado else []
    for _, fila in df.iterrows():
        texto_fila = " | ".join(
            f"{col}: {fila[col]}" for col in df.columns if pd.notna(fila[col])
        )
        if texto_fila:
            lineas.append(texto_fila)
    return "\n".join(lineas)


def extraer_texto_excel(archivo) -> str:
    libro = pd.ExcelFile(archivo)
    partes = []
    for hoja in libro.sheet_names:
        df = libro.parse(hoja)
        partes.append(_dataframe_a_texto(df, encabezado=f"### Hoja: {hoja}"))
    return "\n\n".join(partes)


def extraer_texto_csv(archivo) -> str:
    df = pd.read_csv(archivo)
    return _dataframe_a_texto(df)


def extraer_texto_json(archivo) -> str:
    contenido = json.load(archivo)
    return json.dumps(contenido, ensure_ascii=False, indent=2)


def extraer_texto_html(archivo) -> str:
    soup = BeautifulSoup(archivo.read(), "html.parser")
    return soup.get_text(separator="\n")


def extraer_texto_pptx(archivo) -> str:
    presentacion = Presentation(archivo)
    partes = []
    for i, diapositiva in enumerate(presentacion.slides, start=1):
        partes.append(f"### Diapositiva {i}")
        for forma in diapositiva.shapes:
            if forma.has_text_frame:
                for parrafo in forma.text_frame.paragraphs:
                    texto = "".join(run.text for run in parrafo.runs)
                    if texto.strip():
                        partes.append(texto)
    return "\n".join(partes)


EXTRACTORES = {
    ".pdf": extraer_texto_pdf,
    ".docx": extraer_texto_docx,
    ".txt": extraer_texto_txt,
    ".md": extraer_texto_txt,
    ".xlsx": extraer_texto_excel,
    ".xls": extraer_texto_excel,
    ".csv": extraer_texto_csv,
    ".json": extraer_texto_json,
    ".html": extraer_texto_html,
    ".htm": extraer_texto_html,
    ".pptx": extraer_texto_pptx,
}


def extraer_texto(archivo_subido) -> str:
    _, extension = os.path.splitext(archivo_subido.name.lower())
    extractor = EXTRACTORES.get(extension)
    if not extractor:
        st.warning(f"Formato no soportado: {archivo_subido.name}")
        return ""
    try:
        return extractor(archivo_subido)
    except Exception as error:
        st.error(f"No se pudo procesar {archivo_subido.name}: {error}")
        return ""


# --------------------------------------------------------------------------
# División en fragmentos (chunking) - formato-agnóstico
# --------------------------------------------------------------------------
def dividir_en_fragmentos(texto: str, nombre_archivo: str) -> list[dict]:
    texto = " ".join(texto.split())  # normaliza espacios/saltos de línea
    fragmentos = []
    inicio = 0
    while inicio < len(texto):
        fin = min(inicio + CHUNK_SIZE, len(texto))
        fragmento = texto[inicio:fin]
        if fragmento.strip():
            fragmentos.append({"texto": fragmento, "fuente": nombre_archivo})
        if fin == len(texto):
            break
        inicio = fin - CHUNK_OVERLAP
    return fragmentos


# --------------------------------------------------------------------------
# Modelo de embeddings local (cacheado) + índice FAISS
# --------------------------------------------------------------------------
@st.cache_resource(show_spinner="Cargando modelo de embeddings...")
def cargar_modelo_embeddings():
    from sentence_transformers import SentenceTransformer
    return SentenceTransformer(EMBED_MODEL_NAME)


def construir_indice(fragmentos: list[dict]):
    modelo = cargar_modelo_embeddings()
    textos = [f["texto"] for f in fragmentos]
    embeddings = modelo.encode(textos, show_progress_bar=False, normalize_embeddings=True)
    embeddings = np.array(embeddings, dtype="float32")
    indice = faiss.IndexFlatIP(embeddings.shape[1])  # producto interno = similitud coseno (vectores normalizados)
    indice.add(embeddings)
    return indice


def recuperar_fragmentos_relevantes(pregunta: str, indice, fragmentos: list[dict], k: int = TOP_K):
    modelo = cargar_modelo_embeddings()
    vector_pregunta = modelo.encode([pregunta], normalize_embeddings=True).astype("float32")
    _, posiciones = indice.search(vector_pregunta, min(k, len(fragmentos)))
    return [fragmentos[i] for i in posiciones[0] if i != -1]


# --------------------------------------------------------------------------
# Cliente de la API de Gemini (capa gratuita)
# --------------------------------------------------------------------------
def obtener_api_key() -> str | None:
    if os.environ.get("GEMINI_API_KEY"):
        return os.environ["GEMINI_API_KEY"]
    try:
        return st.secrets.get("GEMINI_API_KEY")
    except Exception:
        return None


def generar_respuesta(pregunta: str, contexto: list[dict], modelo: str, api_key: str) -> str:
    cliente = genai.Client(api_key=api_key)
    bloques_contexto = "\n\n".join(f"[{f['fuente']}]\n{f['texto']}" for f in contexto)
    prompt_usuario = (
        f"CONTEXTO:\n{bloques_contexto}\n\n"
        f"PREGUNTA: {pregunta}\n\n"
        "Responde solo con base en el CONTEXTO anterior."
    )
    respuesta = cliente.models.generate_content(
        model=modelo,
        contents=prompt_usuario,
        config=genai_types.GenerateContentConfig(system_instruction=SYSTEM_PROMPT),
    )
    return respuesta.text


# --------------------------------------------------------------------------
# Estado de sesión
# --------------------------------------------------------------------------
for clave, valor in {
    "fragmentos": [],
    "indice": None,
    "mensajes": [],
    "documentos_cargados": [],
    "api_key_temporal": "",
}.items():
    if clave not in st.session_state:
        st.session_state[clave] = valor

# --------------------------------------------------------------------------
# Barra lateral: carga de documentos y configuración
# --------------------------------------------------------------------------
with st.sidebar:
    st.header("📁 Documentos internos")
    with st.expander("Categorías de referencia sugeridas"):
        for categoria in CATEGORIAS_REFERENCIA:
            st.caption(f"• {categoria}")

    archivos = st.file_uploader(
        "Cargar documentos (PDF, Word, Excel, PowerPoint, Markdown, CSV, JSON, HTML)",
        type=["pdf", "docx", "xlsx", "xls", "pptx", "md", "csv", "json", "html", "htm", "txt"],
        accept_multiple_files=True,
    )

    if st.button("Procesar documentos", type="primary", disabled=not archivos):
        todos_fragmentos = []
        with st.spinner("Extrayendo texto y generando embeddings..."):
            for archivo in archivos:
                texto = extraer_texto(archivo)
                if texto:
                    todos_fragmentos.extend(dividir_en_fragmentos(texto, archivo.name))
            if todos_fragmentos:
                st.session_state.fragmentos = todos_fragmentos
                st.session_state.indice = construir_indice(todos_fragmentos)
                st.session_state.documentos_cargados = [a.name for a in archivos]
                st.session_state.mensajes = []
                st.success(f"{len(archivos)} documento(s) procesados en {len(todos_fragmentos)} fragmentos.")
            else:
                st.error("No se pudo extraer texto de los documentos cargados.")

    if st.session_state.documentos_cargados:
        st.caption("Documentos activos:")
        for nombre in st.session_state.documentos_cargados:
            st.caption(f"• {nombre}")

    if st.button("🗑️ Vaciar documentos"):
        st.session_state.fragmentos = []
        st.session_state.indice = None
        st.session_state.mensajes = []
        st.session_state.documentos_cargados = []
        st.rerun()

    st.divider()
    st.header("⚙️ Modelo (API de Gemini)")
    modelo_seleccionado = st.selectbox("Modelo a usar", MODELOS_GEMINI)

    api_key = obtener_api_key()
    if not api_key:
        st.warning(
            "No se detectó GEMINI_API_KEY en variables de entorno ni en "
            "`.streamlit/secrets.toml`. Puedes ingresarla aquí solo para "
            "esta sesión (no se guarda)."
        )
        st.session_state.api_key_temporal = st.text_input(
            "API key de Gemini", type="password", value=st.session_state.api_key_temporal
        )
        api_key = st.session_state.api_key_temporal or None

    top_k = st.slider("Fragmentos a recuperar por pregunta", 2, 8, TOP_K)

# --------------------------------------------------------------------------
# Área principal: chat
# --------------------------------------------------------------------------
st.title("🏢 Agente Corporativo de IA")
st.caption(
    "Responde únicamente con base en los documentos internos cargados. "
    "Acceso abierto para todos los colaboradores."
)

if not st.session_state.indice:
    st.info("Carga y procesa al menos un documento desde la barra lateral para comenzar.")
elif not api_key:
    st.info("Configura la API key de Gemini en la barra lateral para poder generar respuestas.")

for mensaje in st.session_state.mensajes:
    with st.chat_message(mensaje["role"]):
        st.markdown(mensaje["content"])

pregunta = st.chat_input(
    "Escribe tu pregunta sobre los documentos internos...",
    disabled=not (st.session_state.indice and api_key),
)

if pregunta:
    st.session_state.mensajes.append({"role": "user", "content": pregunta})
    with st.chat_message("user"):
        st.markdown(pregunta)

    with st.chat_message("assistant"):
        with st.spinner("Buscando en los documentos y generando respuesta..."):
            fragmentos_relevantes = recuperar_fragmentos_relevantes(
                pregunta, st.session_state.indice, st.session_state.fragmentos, k=top_k
            )
            try:
                respuesta = generar_respuesta(pregunta, fragmentos_relevantes, modelo_seleccionado, api_key)
            except Exception as error:
                respuesta = f"Ocurrió un error al consultar la API de Gemini: {error}"
            st.markdown(respuesta)
            with st.expander("Fragmentos usados como contexto"):
                for f in fragmentos_relevantes:
                    st.caption(f"**{f['fuente']}**")
                    st.text(f["texto"][:400] + ("..." if len(f["texto"]) > 400 else ""))

    st.session_state.mensajes.append({"role": "assistant", "content": respuesta})
